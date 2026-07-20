import collections
import collections.abc
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # World-Class Corporate Color Palette (Slate, Royal Blue, Soft Cyan, Deep Violet)
    COLOR_PRIMARY_TEXT = RGBColor(15, 23, 42)    # #0F172A (Slate 900 - Premium Dark Title Text)
    COLOR_MUTED_TEXT = RGBColor(71, 85, 105)      # #475569 (Slate 600 - Muted Body Text)
    COLOR_BRAND_BLUE = RGBColor(37, 99, 235)      # #2563EB (Royal Blue - Main Theme Action Color)
    COLOR_CYAN = RGBColor(6, 182, 212)            # #06B6D4 (Teal/Cyan - Primary Accent)
    COLOR_PURPLE = RGBColor(124, 58, 237)         # #7C3AED (Deep Purple/Violet - Secondary Accent)
    COLOR_SHADOW = RGBColor(241, 245, 249)         # #F1F5F9 (Soft Card Shadow Backdrop)
    COLOR_BG = RGBColor(255, 255, 255)            # #FFFFFF (Pure White Base)
    COLOR_BG_LIGHT = RGBColor(248, 250, 252)      # #F8FAFC (Slate 50 background tint)
    COLOR_BORDER_GLASS = RGBColor(226, 232, 240)  # #E2E8F0 (Subtle borders)
    
    # Helper to apply clean background with constellation lines
    def set_website_theme_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG_LIGHT
        
        # Soft blue/cyan gradient meshes (blobs) in corners
        tr_blob = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(-1.5), Inches(4.5), Inches(4.5))
        tr_blob.fill.solid()
        tr_blob.fill.fore_color.rgb = RGBColor(236, 254, 255)  # Softest Cyan
        tr_blob.line.fill.background()
        
        bl_blob = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(5.0), Inches(4.0), Inches(4.0))
        bl_blob.fill.solid()
        bl_blob.fill.fore_color.rgb = RGBColor(245, 243, 255)  # Softest Purple
        bl_blob.line.fill.background()
        
        # Constellation network nodes
        points = [
            (0.8, 1.2), (2.2, 2.5), (1.5, 4.8), (3.2, 6.2), (4.5, 3.5),
            (5.8, 1.5), (7.2, 5.2), (8.5, 2.2), (9.8, 6.2), (11.2, 3.8),
            (12.2, 1.8), (12.5, 6.0), (8.2, 6.5)
        ]
        
        for i, (x1, y1) in enumerate(points):
            for j in range(i + 1, len(points)):
                x2, y2 = points[j]
                dist = ((x1 - x2)**2 + (y1 - y2)**2)**0.5
                if dist < 3.0:
                    try:
                        conn = slide.shapes.add_connector(
                            MSO_CONNECTOR.STRAIGHT, 
                            Inches(x1), Inches(y1), 
                            Inches(x2), Inches(y2)
                        )
                        conn.line.color.rgb = RGBColor(226, 232, 240)
                        conn.line.width = Pt(0.5)
                    except:
                        pass
                        
        for x, y in points:
            node = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - 0.04), Inches(y - 0.04), Inches(0.08), Inches(0.08))
            node.fill.solid()
            node.fill.fore_color.rgb = COLOR_CYAN
            node.line.fill.background()
            
    def set_bg(slide, light_blue=False):
        set_website_theme_bg(slide)
        
    # Helper to draw a glassmorphism card with a soft shadow backdrop
    def add_glass_card(slide, left, top, w, h):
        # 1. Shadow card (slightly offset, grey, no border)
        shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.06), top + Inches(0.06), w, h)
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = COLOR_SHADOW
        shadow.line.fill.background()
        
        # 2. Main card (white, thin border)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_BG
        card.line.color.rgb = COLOR_BORDER_GLASS
        card.line.width = Pt(1.0)
        return card

    # Helper to add standard title with separating line
    def add_storefront_header(slide, title_text):
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Arial"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY_TEXT
        
        # Elegant blue separator line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_BORDER_GLASS
        line.line.fill.background()
        
    def add_title(slide, text):
        add_storefront_header(slide, text)

    # Vector Icon drawing helpers
    def draw_user_icon(slide, left, top):
        head = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.1), top, Inches(0.2), Inches(0.2))
        head.fill.solid()
        head.fill.fore_color.rgb = COLOR_CYAN
        head.line.fill.background()
        
        body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top + Inches(0.22), Inches(0.4), Inches(0.18))
        body.fill.solid()
        body.fill.fore_color.rgb = COLOR_BRAND_BLUE
        body.line.fill.background()

    def draw_search_icon(slide, left, top):
        lens = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(0.24), Inches(0.24))
        lens.fill.solid()
        lens.fill.fore_color.rgb = COLOR_BG_LIGHT
        lens.line.color.rgb = COLOR_BRAND_BLUE
        lens.line.width = Pt(2.0)
        
        handle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.18), top + Inches(0.18), Inches(0.16), Inches(0.08))
        handle.fill.solid()
        handle.fill.fore_color.rgb = COLOR_BRAND_BLUE
        handle.line.fill.background()
        handle.rotation = 45.0

    def draw_db_icon(slide, left, top):
        for i in range(3):
            cyl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top + Inches(i * 0.13), Inches(0.38), Inches(0.1))
            cyl.fill.solid()
            cyl.fill.fore_color.rgb = COLOR_BRAND_BLUE
            cyl.line.color.rgb = COLOR_BORDER_GLASS
            cyl.line.width = Pt(0.5)

    def draw_engine_icon(slide, left, top):
        gear1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top + Inches(0.05), Inches(0.25), Inches(0.25))
        gear1.fill.solid()
        gear1.fill.fore_color.rgb = COLOR_BRAND_BLUE
        gear1.line.fill.background()
        
        gear2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.15), top + Inches(0.15), Inches(0.2), Inches(0.2))
        gear2.fill.solid()
        gear2.fill.fore_color.rgb = COLOR_PURPLE
        gear2.line.fill.background()

    # ==========================================
    # SLIDE 1: Title Slide (Hero Card replica)
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_website_theme_bg(slide)
    
    # Title shadow backing
    sh_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.86), Inches(1.26), Inches(11.7), Inches(5.1))
    sh_card.fill.solid()
    sh_card.fill.fore_color.rgb = COLOR_SHADOW
    sh_card.line.fill.background()
    
    # Storefront Hero Banner Block (Royal Blue)
    hero = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.2), Inches(11.7), Inches(5.1))
    hero.fill.solid()
    hero.fill.fore_color.rgb = COLOR_PRIMARY_TEXT
    hero.line.color.rgb = COLOR_BRAND_BLUE
    
    htf = hero.text_frame
    htf.word_wrap = True
    
    p0 = htf.paragraphs[0]
    p0.text = "AURORA AI RECS ENGINE v2.0"
    p0.font.name = "Arial"
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_CYAN
    p0.space_after = Pt(24)
    
    p1 = htf.add_paragraph()
    p1.text = "Aurora AI Recommendation Engine"
    p1.font.name = "Arial"
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_BG
    p1.space_after = Pt(10)
    
    p2 = htf.add_paragraph()
    p2.text = "Hybrid Product Recommendation System"
    p2.font.name = "Arial"
    p2.font.size = Pt(24)
    p2.font.color.rgb = COLOR_MUTED_TEXT
    p2.space_after = Pt(36)
    
    p3 = htf.add_paragraph()
    p3.text = "Presented by: Krish Verma | MCA"
    p3.font.name = "Arial"
    p3.font.size = Pt(18)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_BG
    
    slide.notes_slide.notes_text_frame.text = (
        "Speaking Notes:\n"
        "1. Introduce yourself: 'Hello everyone, my name is Krish Verma, and I am presenting my MCA final project.'\n"
        "2. State the topic: 'My project is a Hybrid Product Recommendation System, which is an AI/ML based engine.'\n"
        "3. Goal: 'This system recommends products dynamically to shoppers in real-time.'"
    )

    # ==========================================
    # SLIDE 2: Introduction
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_website_theme_bg(slide)
    add_storefront_header(slide, "Introduction")
    
    # Frosted Glassmorphic card
    add_glass_card(slide, Inches(0.8), Inches(1.8), Inches(6.0), Inches(4.5))
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(6.0), Inches(4.5))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_BG
    card.line.color.rgb = COLOR_BORDER_GLASS
    
    tf_c = card.text_frame
    tf_c.word_wrap = True
    p_def = tf_c.paragraphs[0]
    p_def.text = "What is a Recommendation System?"
    p_def.font.name = "Arial"
    p_def.font.size = Pt(20)
    p_def.font.bold = True
    p_def.font.color.rgb = COLOR_BRAND_BLUE
    p_def.space_after = Pt(12)
    
    p_desc = tf_c.add_paragraph()
    p_desc.text = (
        "• A recommendation system predicts what products a customer might like and buy.\n\n"
        "• It analyzes user interests, search history, and reviews to make smart suggestions.\n\n"
        "• It acts as a virtual shopping assistant to make browsing fast and convenient."
    )
    p_desc.font.name = "Calibri"
    p_desc.font.size = Pt(16)
    p_desc.font.color.rgb = COLOR_PRIMARY_TEXT
    p_desc.line_spacing = 1.3
    
    # Right column simple visual highlights
    box_r = slide.shapes.add_textbox(Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.5))
    tf_r = box_r.text_frame
    tf_r.word_wrap = True
    p_ex = tf_r.paragraphs[0]
    p_ex.text = "Key System Goals:"
    p_ex.font.name = "Arial"
    p_ex.font.size = Pt(20)
    p_ex.font.bold = True
    p_ex.font.color.rgb = COLOR_PRIMARY_TEXT
    p_ex.space_after = Pt(18)
    
    goals = [
        ("🎯 Personalization", "Tailor store view dynamically"),
        ("⚡ High Performance", "Load recommended feed in ms"),
        ("📈 Store Engagement", "Boost store conversions and sales")
    ]
    for idx, (title, desc) in enumerate(goals):
        left, top = Inches(7.2), Inches(2.5 + idx * 1.3)
        add_glass_card(slide, left, top, Inches(5.3), Inches(1.0))
        g_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.3), Inches(1.0))
        g_box.fill.solid()
        g_box.fill.fore_color.rgb = COLOR_BG
        g_box.line.color.rgb = COLOR_BORDER_GLASS
        
        gtf = g_box.text_frame
        gtf.word_wrap = True
        gp1 = gtf.paragraphs[0]
        gp1.text = title
        gp1.font.name = "Arial"
        gp1.font.size = Pt(16)
        gp1.font.bold = True
        gp1.font.color.rgb = COLOR_BRAND_BLUE
        gp1.space_after = Pt(2)
        
        gp2 = gtf.add_paragraph()
        gp2.text = desc
        gp2.font.name = "Calibri"
        gp2.font.size = Pt(12)
        gp2.font.color.rgb = COLOR_MUTED_TEXT

    slide.notes_slide.notes_text_frame.text = (
        "Speaking Notes:\n"
        "1. Define the system: 'A recommendation system is an information filter that predicts what a user wants.'\n"
        "2. Explain why it is important: 'It helps users discover products quickly and boosts store sales.'\n"
        "3. Reference real-world cases: 'Platforms like Amazon, Netflix, and Spotify rely heavily on these systems.'"
    )

    # ==========================================
    # SLIDE 3: Problem Statement
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_website_theme_bg(slide)
    add_storefront_header(slide, "Problem Statement")
    
    # Left warning graphic card
    add_glass_card(slide, Inches(0.8), Inches(2.0), Inches(4.5), Inches(4.2))
    warn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(4.5), Inches(4.2))
    warn.fill.solid()
    warn.fill.fore_color.rgb = COLOR_BG_LIGHT
    warn.line.color.rgb = COLOR_BORDER_GLASS
    
    wtf = warn.text_frame
    wtf.word_wrap = True
    wp = wtf.paragraphs[0]
    wp.text = "⚠️ The E-Commerce Challenge\n\nManually searching through massive catalogs is frustrating and drives customers away."
    wp.font.name = "Arial"
    wp.font.size = Pt(20)
    wp.font.bold = True
    wp.font.color.rgb = COLOR_BRAND_BLUE
    wp.alignment = PP_ALIGN.CENTER
    
    # Right column simple bullet points
    box = slide.shapes.add_textbox(Inches(5.8), Inches(2.0), Inches(6.7), Inches(4.2))
    tf = box.text_frame
    tf.word_wrap = True
    
    problems = [
        "Decision Fatigue: Too many choices confuse users.",
        "Catalog Overload: Thousands of products are added daily.",
        "No Personalization: Traditional searches lack personal touch.",
        "Time Consuming: Finding relevant products takes way too long."
    ]
    for idx, prob in enumerate(problems):
        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
        p.text = "• " + prob
        p.font.name = "Calibri"
        p.font.size = Pt(20)
        p.font.color.rgb = COLOR_PRIMARY_TEXT
        p.space_after = Pt(20)
        
    slide.notes_slide.notes_text_frame.text = (
        "Speaking Notes:\n"
        "1. Highlight the core issue: 'Customers face choice overload. There are too many products online.'\n"
        "2. State the consequence: 'Manual search is slow and frustrating. Customers leave if they get confused.'\n"
        "3. Emphasize the solution: 'A personalized recommender is necessary to guide users automatically.'"
    )

    # ==========================================
    # SLIDE 4: Project Objectives
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_website_theme_bg(slide)
    add_storefront_header(slide, "Project Objectives")
    
    # Goal callout card
    sh_goal = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.86), Inches(1.86), Inches(11.7), Inches(1.2))
    sh_goal.fill.solid()
    sh_goal.fill.fore_color.rgb = COLOR_SHADOW
    sh_goal.line.fill.background()
    
    goal = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.2))
    goal.fill.solid()
    goal.fill.fore_color.rgb = COLOR_BRAND_BLUE
    goal.line.fill.background()
    
    gtf = goal.text_frame
    gtf.word_wrap = True
    gp = gtf.paragraphs[0]
    gp.text = "Core Objective:\n\"The main objective is to recommend products based on user interests using Machine Learning.\""
    gp.font.name = "Arial"
    gp.font.size = Pt(20)
    gp.font.bold = True
    gp.font.color.rgb = RGBColor(255, 255, 255)
    gp.alignment = PP_ALIGN.CENTER
    
    # 4 clean sub-objective pillars
    pillars = [
        ("Better User Experience", "Providing a clean and beautiful shopping storefront interface."),
        ("Faster Product Discovery", "Helping shoppers find relevant items in seconds without manual browsing."),
        ("Personalized Results", "Using intelligent ML algorithms to target individual customer tastes."),
        ("Improved Shopping", "Boosting customer satisfaction and engagement across the storefront.")
    ]
    
    for idx, (title, desc) in enumerate(pillars):
        left = Inches(0.8 + (idx % 2) * 6.0)
        top = Inches(3.4 + (idx // 2) * 1.8)
        
        box = slide.shapes.add_textbox(left, top, Inches(5.5), Inches(1.5))
        tf = box.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = "✓ " + title
        p1.font.name = "Arial"
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_BRAND_BLUE
        p1.space_after = Pt(4)
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Calibri"
        p2.font.size = Pt(16)
        p2.font.color.rgb = COLOR_MUTED_TEXT
        
    slide.notes_slide.notes_text_frame.text = (
        "Speaking Notes:\n"
        "1. State the core goal: 'Our main goal is to suggest products based on what the user actually wants.'\n"
        "2. Break down the pillars: 'We achieve this by providing a fast, personalized, and user-friendly interface.'\n"
        "3. Stress that we focus on simple machine learning loops for scalability."
    )

    # ==========================================
    # SLIDE 5: Why Recommendation System?
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_website_theme_bg(slide)
    add_storefront_header(slide, "Why Recommendation System?")
    
    # Left column descriptive box
    box_l = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.5))
    tf_l = box_l.text_frame
    tf_l.word_wrap = True
    
    p_why = tf_l.paragraphs[0]
    p_why.text = "The Business Value of Personalization"
    p_why.font.name = "Arial"
    p_why.font.size = Pt(20)
    p_why.font.bold = True
    p_why.font.color.rgb = COLOR_BRAND_BLUE
    p_why.space_after = Pt(10)
    
    p_why_body = tf_l.add_paragraph()
    p_why_body.text = (
        "• Recommenders drive customer retention and discovery.\n\n"
        "• Up to 35% of Amazon purchases come from recommendation feeds.\n\n"
        "• Fuses search inputs, user clicks, and ratings into cohesive pipelines."
    )
    p_why_body.font.name = "Calibri"
    p_why_body.font.size = Pt(16)
    p_why_body.font.color.rgb = COLOR_PRIMARY_TEXT
    p_why_body.line_spacing = 1.3
    
    # Right column 4 brand cases
    box_r = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.5))
    tf_r = box_r.text_frame
    tf_r.word_wrap = True
    p_r_title = tf_r.paragraphs[0]
    p_r_title.text = "Industry Applications"
    p_r_title.font.name = "Arial"
    p_r_title.font.size = Pt(20)
    p_r_title.font.bold = True
    p_r_title.font.color.rgb = COLOR_PRIMARY_TEXT
    p_r_title.space_after = Pt(18)
    
    brands = [
        ("🛒 Amazon", "Recommended widgets on feed"),
        ("🛍️ Flipkart", "Catalog product matches"),
        ("🎬 Netflix", "Personalized movie feeds"),
        ("🎵 Spotify", "Dynamic playlist discovery")
    ]
    for idx, (bname, desc) in enumerate(brands):
        left_offset = Inches(6.8 + (idx % 2) * 2.95)
        top_offset = Inches(2.7 + (idx // 2) * 1.8)
        
        add_glass_card(slide, left_offset, top_offset, Inches(2.75), Inches(1.4))
        capsule = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_offset, top_offset, Inches(2.75), Inches(1.4))
        capsule.fill.solid()
        capsule.fill.fore_color.rgb = COLOR_BG
        capsule.line.color.rgb = COLOR_BORDER_GLASS
        
        ctf = capsule.text_frame
        ctf.word_wrap = True
        cp1 = ctf.paragraphs[0]
        cp1.text = bname
        cp1.font.name = "Arial"
        cp1.font.size = Pt(16)
        cp1.font.bold = True
        cp1.font.color.rgb = COLOR_BRAND_BLUE
        cp1.space_after = Pt(2)
        
        cp2 = ctf.add_paragraph()
        cp2.text = desc
        cp2.font.name = "Calibri"
        cp2.font.size = Pt(12)
        cp2.font.color.rgb = COLOR_MUTED_TEXT

    slide.notes_slide.notes_text_frame.text = (
        "Speaking Notes:\n"
        "1. Highlight business value: 'Why do we need this? It keeps customers engaged on store catalogs.'\n"
        "2. Cite metrics: 'Over 35% of Amazon's store traffic is generated through recommendation rails.'\n"
        "3. Conclude: 'Every major media or retail platform utilizes a customized recommender.'"
    )

    # ==========================================
    # SLIDE 6: Technologies Used
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_website_theme_bg(slide)
    add_storefront_header(slide, "Technologies Used")
    
    techs = [
        ("🐍 Python", "Core Programming"),
        ("⚡ Flask", "Web Controller"),
        ("🎨 Tailwind CSS", "Store Styling"),
        ("🖥️ HTML & JS", "Client Interface"),
        ("🗄️ SQLite", "Database Layer"),
        ("🐼 Pandas", "Data Structures"),
        ("🔢 NumPy", "Matrix Math"),
        ("🧠 Scikit-Learn", "Content similarity"),
        ("✨ Surprise Lib", "SVD Collaborative"),
        ("🗂️ GitHub", "Version Control"),
        ("🧪 Jupyter", "Model Prototyping")
    ]
    
    for idx, (title, desc) in enumerate(techs):
        left = Inches(0.8 + (idx % 4) * 2.95)
        top = Inches(2.0 + (idx // 4) * 1.6)
        
        add_glass_card(slide, left, top, Inches(2.75), Inches(1.2))
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.75), Inches(1.2))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_BG
        card.line.color.rgb = COLOR_BORDER_GLASS
        
        tf = card.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_BRAND_BLUE
        p.space_after = Pt(2)
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Calibri"
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLOR_MUTED_TEXT
        p2.alignment = PP_ALIGN.CENTER

    slide.notes_slide.notes_text_frame.text = (
        "Speaking Notes:\n"
        "1. List the tools: 'We use Python for programming, Flask as our lightweight web server, and SQLite for data.'\n"
        "2. Explain the ML stack: 'We use Pandas and NumPy for arrays, Scikit-learn for text similarity, and Surprise for SVD.'\n"
        "3. Emphasize that these tools keep the system lightweight and responsive."
    )

    # ==========================================
    # SLIDE 7: Dataset Overview (Table and stats)
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_website_theme_bg(slide)
    add_storefront_header(slide, "Dataset Overview")
    
    # Left stats block
    box_l = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.5))
    tf_l = box_l.text_frame
    tf_l.word_wrap = True
    
    p_title = tf_l.paragraphs[0]
    p_title.text = "Dataset Statistics"
    p_title.font.name = "Arial"
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_BRAND_BLUE
    p_title.space_after = Pt(14)
    
    stats = [
        "👥 Total Simulated Users: 10 active profiles",
        "📦 Total Catalog Products: 500 item records",
        "⭐ Total Rating Records: 1,500 ratings logs",
        "🏷️ Catalog Categories: 4 major categories",
        "📊 User Ratings: Scaled from 1 to 5 stars"
    ]
    for idx, stat in enumerate(stats):
        p = tf_l.add_paragraph()
        p.text = stat
        p.font.name = "Calibri"
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_PRIMARY_TEXT
        p.space_after = Pt(12)
        
    # Right column structural table
    table_shape = slide.shapes.add_table(5, 3, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.2))
    table = table_shape.table
    
    headers = ["Feature", "Data Type", "Description"]
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_BRAND_BLUE
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = "Arial"
            
    rows_data = [
        ("product_id", "Integer", "Primary identifier"),
        ("name", "String", "Name of product"),
        ("category", "String", "Gadget, Style, Home, Fitness"),
        ("description", "String", "Metadata keywords")
    ]
    for r_idx, (feat, dtype, desc) in enumerate(rows_data):
        table.cell(r_idx + 1, 0).text = feat
        table.cell(r_idx + 1, 1).text = dtype
        table.cell(r_idx + 1, 2).text = desc
        
        for col_idx in range(3):
            cell = table.cell(r_idx + 1, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_BG_LIGHT
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = COLOR_PRIMARY_TEXT
                p.font.name = "Calibri"

    slide.notes_slide.notes_text_frame.text = (
        "Speaking Notes:\n"
        "1. Detail the data structure: 'Our schema is structured across Users, Products, and Ratings.'\n"
        "2. Stat summaries: 'We have 500 products split across 4 categories and over 1,500 user rating records.'\n"
        "3. Emphasize indexing: 'This table shows the schema layout for product tracking.'"
    )

    # ==========================================
    # SLIDE 8: System Architecture
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_website_theme_bg(slide)
    add_storefront_header(slide, "System Architecture")
    
    # 7-step flowchart pipeline
    arch_steps = [
        ("User", 0.6, 2.8, 1.4, 1.2),
        ("Website UI", 2.4, 2.8, 1.5, 1.2),
        ("Flask Server", 4.2, 2.8, 1.6, 1.2),
        ("Rec Engine", 6.1, 2.8, 1.6, 1.2),
        ("TF-IDF / SVD", 8.0, 2.8, 1.6, 1.2),
        ("Hybrid Model", 9.9, 2.8, 1.6, 1.2),
        ("Top Products", 11.8, 2.8, 1.4, 1.2)
    ]
    
    for idx, (name, left, top, w, h) in enumerate(arch_steps):
        add_glass_card(slide, Inches(left), Inches(top), Inches(w), Inches(h))
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
        shape.fill.solid()
        if "Engine" in name or "Hybrid" in name:
            shape.fill.fore_color.rgb = COLOR_BRAND_BLUE
            text_color = COLOR_BG
        else:
            shape.fill.fore_color.rgb = COLOR_BG
            text_color = COLOR_PRIMARY_TEXT
            
        shape.line.color.rgb = COLOR_BORDER_GLASS
        
        # Draw icons
        if name == "User":
            draw_user_icon(slide, Inches(left + 0.5), Inches(top - 0.75))
        elif "Website" in name:
            draw_search_icon(slide, Inches(left + 0.55), Inches(top - 0.65))
        elif "Server" in name:
            draw_db_icon(slide, Inches(left + 0.55), Inches(top - 0.65))
        elif "Engine" in name:
            draw_engine_icon(slide, Inches(left + 0.55), Inches(top - 0.65))
            
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.name = "Arial"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER
        
        # Connecting arrows
        if idx < 6:
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(left + w), Inches(3.25), Inches(0.12), Inches(0.12))
            arr.fill.solid()
            arr.fill.fore_color.rgb = COLOR_CYAN
            arr.line.fill.background()
            
    # Bottom description box
    add_glass_card(slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(1.5))
    desc_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.8), Inches(11.7), Inches(1.5))
    desc_box.fill.solid()
    desc_box.fill.fore_color.rgb = COLOR_BG
    desc_box.line.color.rgb = COLOR_BORDER_GLASS
    
    desc_tf = desc_box.text_frame
    desc_tf.word_wrap = True
    dp = desc_tf.paragraphs[0]
    dp.text = "Flow description: The active user interacts with the HTML interface. Request parameters pass to the Flask backend controller, which coordinates the Recommendation Engine. The engine executes parallel checks on TF-IDF cosine descriptions matrix and SVD rating indices, combining scores inside the Hybrid module."
    dp.font.name = "Calibri"
    dp.font.size = Pt(16)
    dp.font.color.rgb = COLOR_MUTED_TEXT

    slide.notes_slide.notes_text_frame.text = (
        "Speaking Notes:\n"
        "1. Explain the architectural flow: 'A user search query triggers a call to the central recommendation engine.'\n"
        "2. Explain the core math: 'The engine coordinates Content similarity and Collaborative rating matrices in parallel.'\n"
        "3. State the output: 'They are blended at the Hybrid Scoring stage to display the Top 5 recommended items.'"
    )

    # ==========================================
    # SLIDE 9: Project Workflow
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_website_theme_bg(slide)
    add_storefront_header(slide, "Project Workflow")
    
    # 6-step workflow infographic
    workflow_steps = [
        ("User", "Shopper visits storefront catalog page."),
        ("Search", "Inputs keyword or sets alpha tuning weight."),
        ("Database", "Flask queries product & rating records SQL."),
        ("ML Model", "Calculates TF-IDF and SVD user ratings."),
        ("Recommendation", "Blends metrics into dynamic hybrid scores."),
        ("Website", "Renders personalized product feed cards.")
    ]
    
    for idx, (title, desc) in enumerate(workflow_steps):
        left = Inches(0.6 + idx * 2.05)
        top = Inches(2.3)
        w = Inches(1.8)
        h = Inches(3.2)
        
        add_glass_card(slide, left, top, w, h)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_BG
        box.line.color.rgb = COLOR_BORDER_GLASS
        
        # Draw icons
        if title == "User":
            draw_user_icon(slide, Inches(left + 0.65), Inches(top + 0.2))
        elif title == "Search":
            draw_search_icon(slide, Inches(left + 0.72), Inches(top + 0.3))
        elif title == "Database":
            draw_db_icon(slide, Inches(left + 0.72), Inches(top + 0.3))
            
        tf = box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"\n\n\nStep {idx+1}\n{title}"
        p.font.name = "Arial"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_BRAND_BLUE
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(4)
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Calibri"
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_MUTED_TEXT
        p2.alignment = PP_ALIGN.CENTER
        
        if idx < 5:
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(left + w), Inches(3.6), Inches(0.2), Inches(0.2))
            arr.fill.solid()
            arr.fill.fore_color.rgb = COLOR_CYAN
            arr.line.fill.background()

    slide.notes_slide.notes_text_frame.text = (
        "Speaking Notes:\n"
        "1. Explain the 4-step workflow: search, database lookup, model execution, and layout rendering.\n"
        "2. Explain that model execution runs instantly on the server side using serialized pickle files.\n"
        "3. Emphasize that the user sees their updated feed in milliseconds."
    )

    # ==========================================
    # SLIDE 10: Machine Learning Models
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_website_theme_bg(slide)
    add_storefront_header(slide, "Machine Learning Models")
    
    col_w = Inches(3.6)
    h = Inches(4.5)
    
    # Card 1: Content-Based
    add_glass_card(slide, Inches(0.8), Inches(1.8), col_w, h)
    c1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), col_w, h)
    c1.fill.solid()
    c1.fill.fore_color.rgb = COLOR_BG
    c1.line.color.rgb = COLOR_BORDER_GLASS
    tf1 = c1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "📝 Content-Based Filtering"
    p1.font.name = "Arial"
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_BRAND_BLUE
    p1.space_after = Pt(12)
    p1_body = tf1.add_paragraph()
    p1_body.text = (
        "• Recommends items based on product descriptions.\n\n"
        "• Uses TF-IDF for text parsing.\n\n"
        "• Uses Cosine Similarity to calculate matches.\n\n"
        "👉 Example:\n"
        "If a user views an iPhone, system recommends similar mobile phones."
    )
    p1_body.font.name = "Calibri"
    p1_body.font.size = Pt(15)
    p1_body.font.color.rgb = COLOR_PRIMARY_TEXT
    
    # Card 2: Collaborative
    add_glass_card(slide, Inches(4.85), Inches(1.8), col_w, h)
    c2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.85), Inches(1.8), col_w, h)
    c2.fill.solid()
    c2.fill.fore_color.rgb = COLOR_BG
    c2.line.color.rgb = COLOR_BORDER_GLASS
    tf2 = c2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "👥 Collaborative Filtering"
    p2.font.name = "Arial"
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_BRAND_BLUE
    p2.space_after = Pt(12)
    p2_body = tf2.add_paragraph()
    p2_body.text = (
        "• Recommends items based on user ratings and behavior.\n\n"
        "• Uses SVD (Matrix Factorization) to predict ratings.\n\n"
        "• Finds hidden customer tastes.\n\n"
        "👉 Example:\n"
        "Users who liked Laptop also liked Laptop Mouse."
    )
    p2_body.font.name = "Calibri"
    p2_body.font.size = Pt(15)
    p2_body.font.color.rgb = COLOR_PRIMARY_TEXT

    # Card 3: Hybrid (Violet highlight)
    sh_c3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.96), Inches(1.86), col_w, h)
    sh_c3.fill.solid()
    sh_c3.fill.fore_color.rgb = COLOR_SHADOW
    sh_c3.line.fill.background()
    
    c3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.9), Inches(1.8), col_w, h)
    c3.fill.solid()
    c3.fill.fore_color.rgb = COLOR_BRAND_BLUE
    c3.line.color.rgb = COLOR_BRAND_BLUE
    tf3 = c3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = "🔄 Hybrid Recommendation"
    p3.font.name = "Arial"
    p3.font.size = Pt(18)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_BG
    p3.space_after = Pt(12)
    p3_body = tf3.add_paragraph()
    p3_body.text = (
        "\"Combines both methods to improve recommendation accuracy.\"\n\n"
        "• Avoids the shortcomings of cold-starts.\n\n"
        "• Presents a balanced set of product results directly onto the shopper's feed."
    )
    p3_body.font.name = "Calibri"
    p3_body.font.size = Pt(16)
    p3_body.font.color.rgb = COLOR_BG

    slide.notes_slide.notes_text_frame.text = (
        "Speaking Notes:\n"
        "1. Explain Content-Based: 'It matches words in product descriptions. iPhone leads to other smartphones.'\n"
        "2. Explain Collaborative: 'It matches users with similar ratings. Laptop buyers matching with a mouse.'\n"
        "3. Conclude with Hybrid: 'Combining both gives us robust and accurate recommendations.'"
    )

    # ==========================================
    # SLIDE 11: Content-Based Filtering
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_website_theme_bg(slide)
    add_storefront_header(slide, "Content-Based Filtering")
    
    # Left column details
    box_l = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.5))
    tf_l = box_l.text_frame
    tf_l.word_wrap = True
    
    p_title = tf_l.paragraphs[0]
    p_title.text = "Metadata & TF-IDF Extraction"
    p_title.font.name = "Arial"
    p_title.font.size = Pt(20)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_BRAND_BLUE
    p_title.space_after = Pt(10)
    
    p_body = tf_l.add_paragraph()
    p_body.text = (
        "• Recommends items based on product tags, names, and descriptions.\n\n"
        "• TF-IDF calculates how unique a word is across the dataset.\n\n"
        "• Cosine Similarity measures the angle between keyword vectors.\n\n"
        "👉 Key advantage: Resolves recommendation logic for products without ratings."
    )
    p_body.font.name = "Calibri"
    p_body.font.size = Pt(16)
    p_body.font.color.rgb = COLOR_PRIMARY_TEXT
    p_body.line_spacing = 1.3
    
    # Right column visual card
    add_glass_card(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.5))
    box_r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.5))
    box_r.fill.solid()
    box_r.fill.fore_color.rgb = COLOR_BG_LIGHT
    box_r.line.color.rgb = COLOR_BORDER_GLASS
    
    tf_r = box_r.text_frame
    tf_r.word_wrap = True
    pr = tf_r.paragraphs[0]
    pr.text = "📱 Recommendation Pipeline Example:\n\nUser views iPhone 13 Pro ➔ Systems parses description ('iOS', 'mobile', '5G') ➔ Calculates cosine similarity matrices ➔ Recommends iPhone 14 Pro, Samsung Galaxy, OnePlus 11."
    pr.font.name = "Calibri"
    pr.font.size = Pt(18)
    pr.font.color.rgb = COLOR_PRIMARY_TEXT
    pr.space_after = Pt(10)

    slide.notes_slide.notes_text_frame.text = (
        "Speaking Notes:\n"
        "1. Explain Content-Based: 'It matches terms in product descriptions.'\n"
        "2. Detail TF-IDF: 'It measures the uniqueness of words to create keyword vectors.'\n"
        "3. Detail Cosine Similarity: 'It computes similarity scores to match items like related phones.'"
    )

    # ==========================================
    # SLIDE 12: Collaborative Filtering
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_website_theme_bg(slide)
    add_storefront_header(slide, "Collaborative Filtering")
    
    # Left column details
    box_l = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.5))
    tf_l = box_l.text_frame
    tf_l.word_wrap = True
    
    p_title = tf_l.paragraphs[0]
    p_title.text = "SVD (Matrix Factorization) Algorithm"
    p_title.font.name = "Arial"
    p_title.font.size = Pt(20)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_BRAND_BLUE
    p_title.space_after = Pt(10)
    
    p_body = tf_l.add_paragraph()
    p_body.text = (
        "• Analyzes ratings history across users to find overlaps.\n\n"
        "• SVD splits User-Item ratings matrix into hidden feature vectors.\n\n"
        "• Predicts user ratings for products they haven't seen.\n\n"
        "👉 Key advantage: Discovers non-obvious cross-category recommendations."
    )
    p_body.font.name = "Calibri"
    p_body.font.size = Pt(16)
    p_body.font.color.rgb = COLOR_PRIMARY_TEXT
    p_body.line_spacing = 1.3
    
    # Right column comparison card
    add_glass_card(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.5))
    box_r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.5))
    box_r.fill.solid()
    box_r.fill.fore_color.rgb = COLOR_BG_LIGHT
    box_r.line.color.rgb = COLOR_BORDER_GLASS
    
    tf_r = box_r.text_frame
    tf_r.word_wrap = True
    pr = tf_r.paragraphs[0]
    pr.text = "👥 Collaborative Example:\n\nUser A likes Laptop, User B likes Laptop and Mouse ➔ System connects User A to User B ➔ Recommends Mouse to User A based on ratings alignment."
    pr.font.name = "Calibri"
    pr.font.size = Pt(18)
    pr.font.color.rgb = COLOR_PRIMARY_TEXT

    slide.notes_slide.notes_text_frame.text = (
        "Speaking Notes:\n"
        "1. Explain Collaborative: 'It uses user ratings matrix alignment.'\n"
        "2. Explain SVD: 'Singular Value Decomposition splits this matrix into low-rank representations.'\n"
        "3. State benefits: 'It predicts ratings for unseen items based on similar user habits.'"
    )

    # ==========================================
    # SLIDE 13: Website Features
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_website_theme_bg(slide)
    add_storefront_header(slide, "Website Features")
    
    pages = [
        ("Home Page storefront", "Landing Hero panel displaying dynamic SVD recommendation carousels."),
        ("Product Details Page", "Split view panel displaying item descriptions and Content-based similar choices."),
        ("History Logs Page", "Active database logger tracking click types to simulate system logs."),
        ("Analytics Dashboard", "Visual dashboard showing rating distributions and correlation matrix plots.")
    ]
    
    for idx, (title, desc) in enumerate(pages):
        left = Inches(0.8 + (idx % 2) * 6.0)
        top = Inches(1.8 + (idx // 2) * 2.6)
        
        # Browser mockup window with soft shadow
        add_glass_card(slide, left, top, Inches(5.7), Inches(2.2))
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.7), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_BG
        card.line.color.rgb = COLOR_BORDER_GLASS
        
        # Header bar
        header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(5.7), Inches(0.45))
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = COLOR_BORDER_GLASS
        header_bar.line.fill.background()
        
        # Dot controls
        for c_idx in range(3):
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.15 + c_idx * 0.2), top + Inches(0.12), Inches(0.12), Inches(0.12))
            circle.fill.solid()
            circle.fill.fore_color.rgb = COLOR_BRAND_BLUE
            circle.line.fill.background()
            
        tf_box = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.55), Inches(5.1), Inches(1.5))
        tf = tf_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_BRAND_BLUE
        p.space_after = Pt(2)
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Calibri"
        p2.font.size = Pt(13)
        p2.font.color.rgb = COLOR_MUTED_TEXT

    slide.notes_slide.notes_text_frame.text = (
        "Speaking Notes:\n"
        "1. Review the interface screens: Home, Product Details, History, and Analytics.\n"
        "2. Explain that the design uses white background containers with blue text indicators for clarity.\n"
        "3. Highlight that the screens are designed to match modern SaaS templates."
    )

    # ==========================================
    # SLIDE 14: Results & Performance
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_website_theme_bg(slide)
    add_storefront_header(slide, "Results & Performance")
    
    # Left column details
    box_l = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
    tf_l = box_l.text_frame
    tf_l.word_wrap = True
    
    p_title = tf_l.paragraphs[0]
    p_title.text = "System Metric Evaluation"
    p_title.font.name = "Arial"
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_BRAND_BLUE
    p_title.space_after = Pt(12)
    
    results = [
        "Fast Recommendations",
        "   - Response time is optimized using pre-calculated weights in pickle files.",
        "Personalized Results",
        "   - Recommendations dynamically shift immediately upon changing user profiles.",
        "Easy to Use & Modern Interface",
        "   - Minimal text fields, synced search inputs, and responsive layout grids."
    ]
    for idx, r in enumerate(results):
        p = tf_l.add_paragraph()
        p.text = r
        if r.strip().startswith("-"):
            p.font.name = "Calibri"
            p.font.size = Pt(16)
            p.font.color.rgb = COLOR_MUTED_TEXT
            p.space_after = Pt(12)
        else:
            p.font.name = "Calibri"
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = COLOR_PRIMARY_TEXT
            p.space_after = Pt(4)
            
    # Right column visual summary box
    sh_box_r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.56), Inches(2.26), Inches(4.8), Inches(3.8))
    sh_box_r.fill.solid()
    sh_box_r.fill.fore_color.rgb = COLOR_SHADOW
    sh_box_r.line.fill.background()
    
    box_r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(2.2), Inches(4.8), Inches(3.8))
    box_r.fill.solid()
    box_r.fill.fore_color.rgb = COLOR_BRAND_BLUE
    box_r.line.color.rgb = COLOR_BRAND_BLUE
    tf_r = box_r.text_frame
    tf_r.word_wrap = True
    pr = tf_r.paragraphs[0]
    pr.text = "⚡ Evaluation Strengths\n\n• Scalable architecture\n\n• High recommendation relevance\n\n• No complex server prerequisites\n\n• Dynamic weight sliders"
    pr.font.name = "Arial"
    pr.font.size = Pt(20)
    pr.font.bold = True
    pr.font.color.rgb = COLOR_BG
    pr.alignment = PP_ALIGN.LEFT
    pr.space_after = Pt(10)

    slide.notes_slide.notes_text_frame.text = (
        "Speaking Notes:\n"
        "1. Explain that recommendations load instantly due to weight caching.\n"
        "2. Highlight that the suggestions are highly personalized and adapt as user tastes change.\n"
        "3. Emphasize that the interface is simple, lightweight, and modern."
    )

    # ==========================================
    # SLIDE 15: Advantages & Limitations
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_website_theme_bg(slide)
    add_storefront_header(slide, "Advantages & Limitations")
    
    # Left Card: Advantages
    add_glass_card(slide, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.5))
    adv = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.5))
    adv.fill.solid()
    adv.fill.fore_color.rgb = COLOR_BG
    adv.line.color.rgb = COLOR_BORDER_GLASS
    tf_a = adv.text_frame
    tf_a.word_wrap = True
    pa = tf_a.paragraphs[0]
    pa.text = "🟢 Key Advantages"
    pa.font.name = "Arial"
    pa.font.size = Pt(20)
    pa.font.bold = True
    pa.font.color.rgb = COLOR_BRAND_BLUE
    pa.space_after = Pt(14)
    pa_body = tf_a.add_paragraph()
    pa_body.text = (
        "• High Accuracy: Combines textual metadata & user ratings.\n\n"
        "• Resolves Cold-Start: Metadata (TF-IDF) helps rank new items.\n\n"
        "• Fast Execution: Cache pre-computation yields speedy loads.\n\n"
        "• Easy Scaling: SQLite structure makes integration modular."
    )
    pa_body.font.name = "Calibri"
    pa_body.font.size = Pt(16)
    pa_body.font.color.rgb = COLOR_PRIMARY_TEXT
    
    # Right Card: Limitations
    add_glass_card(slide, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.5))
    lim = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.5))
    lim.fill.solid()
    lim.fill.fore_color.rgb = COLOR_BG
    lim.line.color.rgb = COLOR_BORDER_GLASS
    tf_l = lim.text_frame
    tf_l.word_wrap = True
    pl = tf_l.paragraphs[0]
    pl.text = "🔴 System Limitations"
    pl.font.name = "Arial"
    pl.font.size = Pt(20)
    pl.font.bold = True
    pl.font.color.rgb = COLOR_PURPLE
    pl.space_after = Pt(14)
    pl_body = tf_l.add_paragraph()
    pl_body.text = (
        "• Rating Sparsity: SVD needs minimum baseline rating inputs.\n\n"
        "• Static Cache: Retraining needed for large, real-time products updates.\n\n"
        "• Computation Load: Large matrices require high RAM bounds.\n\n"
        "• Keyword Reliance: TF-IDF depends on good description quality."
    )
    pl_body.font.name = "Calibri"
    pl_body.font.size = Pt(16)
    pl_body.font.color.rgb = COLOR_PRIMARY_TEXT

    slide.notes_slide.notes_text_frame.text = (
        "Speaking Notes:\n"
        "1. Detail advantages: 'Hybrid scoring blends descriptions and user behavior to bypass cold-start.'\n"
        "2. Address limits: 'Like all rating systems, SVD matrices face sparsity limitations.'\n"
        "3. Emphasize that cache files resolve real-time response delay."
    )

    # ==========================================
    # SLIDE 16: Future Scope
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_website_theme_bg(slide)
    add_storefront_header(slide, "Future Scope")
    
    scopes = [
        ("🤖 AI Chatbot", "Natural language chat assistants for guided shopping recommendations."),
        ("🗣️ Voice Search", "Enable search queries using voice triggers on client browsers."),
        ("👤 User Login", "Creating logins and signups databases to maintain user logs."),
        ("📊 Admin Dashboard", "Analytics dashboard for store owners to trace customer activity."),
        ("🧠 Deep Learning", "Upgrading core algorithms from SVD to Neural Collaborative models."),
        ("📱 Mobile Application", "Building Flutter or React Native mobile shopping catalogs.")
    ]
    
    for idx, (title, desc) in enumerate(scopes):
        left = Inches(0.8 + (idx % 3) * 4.0)
        top = Inches(2.0 + (idx // 3) * 2.3)
        
        add_glass_card(slide, left, top, Inches(3.7), Inches(2.0))
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.7), Inches(2.0))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_BG
        box.line.color.rgb = COLOR_BORDER_GLASS
        
        tf = box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_BRAND_BLUE
        p.space_after = Pt(4)
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Calibri"
        p2.font.size = Pt(14)
        p2.font.color.rgb = COLOR_MUTED_TEXT

    slide.notes_slide.notes_text_frame.text = (
        "Speaking Notes:\n"
        "1. Discuss future plans: adding chat assistants, voice search, and login systems.\n"
        "2. Mention scaling algorithms: migrating to Deep Learning recommendation models.\n"
        "3. Emphasize that the codebase is designed to be easily extensible for these features."
    )

    # ==========================================
    # SLIDE 17: Conclusion
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_website_theme_bg(slide)
    add_storefront_header(slide, "Conclusion")
    
    # Left vertical accent bar
    side_border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(0.1), Inches(4.0))
    side_border.fill.solid()
    side_border.fill.fore_color.rgb = COLOR_BRAND_BLUE
    side_border.line.fill.background()
    
    # Content block
    c_box = slide.shapes.add_textbox(Inches(1.1), Inches(2.2), Inches(11.2), Inches(4.0))
    c_tf = c_box.text_frame
    c_tf.word_wrap = True
    
    cp = c_tf.paragraphs[0]
    cp.text = "Summary & Key Takeaways"
    cp.font.name = "Arial"
    cp.font.size = Pt(24)
    cp.font.bold = True
    cp.font.color.rgb = COLOR_PRIMARY_TEXT
    cp.space_after = Pt(14)
    
    cp2 = c_tf.add_paragraph()
    cp2.text = (
        "\"This project helps users find products quickly using Machine Learning. "
        "It combines Content-Based Filtering and Collaborative Filtering to provide better recommendations. "
        "The system is simple, fast, and user-friendly.\"\n\n"
        "By fusing metadata keyword extraction and user collaborative SVD rating matrices, "
        "the application successfully addresses the cold-start problem and outputs highly relevant results "
        "suited for real-world e-commerce deployment."
    )
    cp2.font.name = "Calibri"
    cp2.font.size = Pt(20)
    cp2.font.color.rgb = COLOR_MUTED_TEXT
    cp2.line_spacing = 1.3

    slide.notes_slide.notes_text_frame.text = (
        "Speaking Notes:\n"
        "1. Sum up: 'Our project demonstrates how Machine Learning improves shopping personalization.'\n"
        "2. Reiterate the hybrid benefit: 'Combining Content and Collaborative filtering gives us the best of both worlds.'\n"
        "3. End with a confidence statement: 'The system is fast, modular, and ready for deployment.'"
    )

    # ==========================================
    # SLIDE 18: Thank You
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_website_theme_bg(slide)
    
    # Shadow for thank you card
    sh_thank = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.56), Inches(1.56), Inches(10.333), Inches(4.5))
    sh_thank.fill.solid()
    sh_thank.fill.fore_color.rgb = COLOR_SHADOW
    sh_thank.line.fill.background()
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_BG
    card.line.color.rgb = COLOR_BORDER_GLASS
    
    c_tf = card.text_frame
    c_tf.word_wrap = True
    
    p = c_tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.name = "Arial"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = COLOR_BRAND_BLUE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(20)
    
    p2 = c_tf.add_paragraph()
    p2.text = "Questions & Answers Session\n\nKrish Verma | MCA"
    p2.font.name = "Arial"
    p2.font.size = Pt(24)
    p2.font.color.rgb = COLOR_PRIMARY_TEXT
    p2.alignment = PP_ALIGN.CENTER
    
    slide.notes_slide.notes_text_frame.text = (
        "Speaking Notes:\n"
        "1. Conclude: 'Thank you everyone for your time.'\n"
        "2. Open the floor: 'I am now happy to answer any questions about the algorithms, dataset, or implementation.'"
    )

    # Save presentation
    output_filename = "Aurora_AI_Recommendation_System_Presentation.pptx"
    prs.save(output_filename)
    print(f"Presentation successfully compiled to {output_filename}!")

if __name__ == "__main__":
    build_presentation()
