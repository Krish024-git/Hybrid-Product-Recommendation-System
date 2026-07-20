import collections
import collections.abc
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set dimensions to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Define custom color palette
    BG_COLOR = RGBColor(248, 251, 255)       # #F8FBFF (Light Blue base)
    PRIMARY_COLOR = RGBColor(37, 99, 235)     # #2563EB (Royal Blue)
    ACCENT_COLOR = RGBColor(59, 130, 246)     # #3B82F6 (Sky Blue)
    TEXT_MAIN = RGBColor(15, 23, 42)          # #0F172A (Slate 900)
    TEXT_MUTED = RGBColor(71, 85, 105)        # #475569 (Slate 600)
    
    # Helper to apply clean background color
    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        
    # Helper to add standard header block
    def add_header(slide, title_text):
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Arial"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        
        # Subtle horizontal accent line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(2.0), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_COLOR
        line.line.fill.background()
        
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    
    # Left accent side border shape
    side_border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    side_border.fill.solid()
    side_border.fill.fore_color.rgb = PRIMARY_COLOR
    side_border.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.0), Inches(4.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = "AURORA AI RECOMMENDATION ENGINE"
    p0.font.name = "Arial"
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_COLOR
    p0.space_after = Pt(20)
    
    p1 = tf.add_paragraph()
    p1.text = "Hybrid Product Recommendation System"
    p1.font.name = "Arial"
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_MAIN
    p1.space_after = Pt(12)
    
    p2 = tf.add_paragraph()
    p2.text = "A real-time personalized e-commerce shopping experience"
    p2.font.name = "Arial"
    p2.font.size = Pt(20)
    p2.font.color.rgb = TEXT_MUTED
    p2.space_after = Pt(36)
    
    p3 = tf.add_paragraph()
    p3.text = "Presented by: Project Team"
    p3.font.name = "Arial"
    p3.font.size = Pt(14)
    p3.font.bold = True
    p3.font.color.rgb = PRIMARY_COLOR
    
    # Helper to generate bullet point content
    def add_bullets_slide(title, items, speaking_note=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide)
        add_header(slide, title)
        
        textbox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
        tf = textbox.text_frame
        tf.word_wrap = True
        
        for idx, item in enumerate(items):
            p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
            p.text = item
            
            # Basic spacing and indentation
            if item.strip().startswith("-") or item.strip().startswith("   "):
                p.font.name = "Calibri"
                p.font.size = Pt(18)
                p.font.color.rgb = TEXT_MUTED
                p.space_after = Pt(10)
            else:
                p.font.name = "Calibri"
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = TEXT_MAIN
                p.space_after = Pt(18)
                
        if speaking_note:
            slide.notes_slide.notes_text_frame.text = speaking_note
            
    # Slide 2: The Problem
    add_bullets_slide(
        "The E-Commerce Problem Statement",
        [
            "Choice Overload in Online Retail",
            "   - Customers are overwhelmed by thousands of products, leading to decision fatigue.",
            "Static Catalog Structures",
            "   - Standard storefronts display the same homepage catalog cards to everyone.",
            "Cold Start Challenge",
            "   - Inability to suggest relevant products to new users with zero purchase history.",
            "Lost Store Conversion",
            "   - Lower user engagement leads to abandoned carts and lost business revenue."
        ],
        "Speaking Notes:\n- Explain that users leave stores when they don't find what they like immediately."
    )
    
    # Slide 3: Our Solution
    add_bullets_slide(
        "The Solution: Aurora AI Engine",
        [
            "Dynamic Personalization Layer",
            "   - Fuses user rating profiles and product metadata for tailored suggestions.",
            "Collaborative & Content Blending",
            "   - Uses multiple algorithms in parallel to overcome individual limitations.",
            "Premium Client Visuals",
            "   - Modern light theme, animated constellation background, and interactive interface.",
            "Localized Assets",
            "   - Integrated Indian Rupee (₹) catalog and realistic product category photography."
        ],
        "Speaking Notes:\n- Emphasize how the hybrid approach solves the cold start and increases CTR."
    )
    
    # Slide 4: Algorithmic Architecture
    add_bullets_slide(
        "Recommendation System Architecture",
        [
            "Collaborative Filtering (SVD Matrix Factorization)",
            "   - Predicts rating scores for unviewed items using historical review matrices.",
            "Content-Based Filtering (TF-IDF & Cosine Similarity)",
            "   - Analyzes description keywords to list matching alternative catalog choices.",
            "Hybrid Score Combiner",
            "   - Blends item feature matches and rating predictions for maximum relevance."
        ],
        "Speaking Notes:\n- Describe how SVD reads user-rating matrices and TF-IDF checks product descriptions."
    )
    
    # Slide 5: System Technology Stack
    add_bullets_slide(
        "System Technology Stack",
        [
            "Backend Web Controller",
            "   - Python with Flask micro-framework handling routing API requests.",
            "Database Engine",
            "   - Relational SQLite database managing product attributes, users, and logs.",
            "Data & Machine Learning Core",
            "   - Sci-kit Learn (TF-IDF Vectorizer), Pandas dataframes, and Custom SVD solver.",
            "Responsive Client Interface",
            "   - HTML5, Tailwind CSS layout utilities, and interactive background JS Canvas."
        ],
        "Speaking Notes:\n- Highlight that SQLite and Flask are extremely easy to scale and deploy."
    )
    
    # Slide 6: Key Live Demo Highlights
    add_bullets_slide(
        "Key Live Demo Highlights",
        [
            "Active Simulated Profile Dropdown",
            "   - Swap active buyer profile dynamically to trigger immediate recommendation shifts.",
            "Real-Time Synchronized Search",
            "   - Interactive double-search inputs filtering catalog items instantly on-key.",
            "Interaction History Table",
            "   - Full logging trace mapping click triggers and prediction matching algorithms.",
            "Integrated Model Analytics",
            "   - Clean dashboard showing distribution curves and correlation heatmaps."
        ],
        "Speaking Notes:\n- Point to the active user profile switcher during your live demo."
    )
    
    # Slide 7: Future Scope
    add_bullets_slide(
        "Future Enhancements",
        [
            "Deep Learning Matrix Factorization",
            "   - Upgrading to Neural Collaborative Filtering (NCF) using PyTorch.",
            "Integrated Payment Gateways",
            "   - Adding standard checkout hooks (UPI, Cards) to complete order flows.",
            "Live A/B Model Testing Panel",
            "   - Dynamic dashboard monitoring user clicks to optimize model weights live."
        ],
        "Speaking Notes:\n- Conclude by showing the path to commercializing the project."
    )
    
    # Slide 8: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.3), Inches(3.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.name = "Arial"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "Questions & Answers Session"
    p2.font.name = "Arial"
    p2.font.size = Pt(24)
    p2.font.color.rgb = TEXT_MAIN
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)
    
    output_filename = "Aurora_AI_Recommendation_Presentation.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully as {output_filename}!")

if __name__ == "__main__":
    create_presentation()
